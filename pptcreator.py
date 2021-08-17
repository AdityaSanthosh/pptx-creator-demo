from pptx import Presentation

# PRESENTATION OBJECT
prs = Presentation()

# LAYOUT AND SLIDES USED IN THIS PRESENTATION
SLIDE_LAYOUT_TITLE_SLIDE = 0
SLIDE_LAYOUT_PICTURE_WITH_CAPTION = 8

title_slide_layout = prs.slide_layouts[SLIDE_LAYOUT_TITLE_SLIDE]
image_slide_layout = prs.slide_layouts[SLIDE_LAYOUT_PICTURE_WITH_CAPTION]


# TITLE SLIDE
title_slide = prs.slides.add_slide(title_slide_layout)
title_slide_title = title_slide.shapes.title
title_slide_subtitle = title_slide.placeholders[1]

title_slide_title.text = "Nike Products"
title_slide_subtitle.text = "These are the products of Nike"


slidetitles = ["Home Garden", "Orange Juice!", "Simple Work Place", "Modern Furniture", "Cameraaa!!"]
slidecaptions = ["Plants in a house", "Tasty!", "Aesthetically Awesome ", "Too Expensive", "Analog Awesome"]

# Generating Slides Programmatically on a loop
for i in range(1, 6):
    eachslide = prs.slides.add_slide(image_slide_layout)
    eachslide_title = eachslide.placeholders[0]
    eachslide_pic = eachslide.placeholders[1].insert_picture("image_"+str(i)+'.png')
    eachslide_caption = eachslide.placeholders[2]
    eachslide_title.text = slidetitles[i-1] #the index is i-1 because imagefiles start from 0 and slide lists from 0
    eachslide_caption.text = slidecaptions[i-1]

"""
Another Implementation
# IMAGE1 SLIDE
Image1_slide = prs.slides.add_slide(image_slide_layout)
Image1_slide_title = Image1_slide.placeholders[0]
pic1 = Image1_slide.placeholders[1].insert_picture('image1.jpg')
Image1_slide_caption = Image1_slide.placeholders[2]

Image1_slide_title.text = "Home Garden"
Image1_slide_caption.text = "Plants in a house"

# IMAGE2 SLIDE
Image2_slide = prs.slides.add_slide(image_slide_layout)
Image2_slide_title = Image2_slide.placeholders[0]
pic2 = Image2_slide.placeholders[1].insert_picture('image2.jpg')
Image2_slide_caption = Image2_slide.placeholders[2]

Image2_slide_title.text = "Orange Guice!"
Image2_slide_caption.text = "It's good"

# IMAGE3 SLIDE
Image3_slide = prs.slides.add_slide(image_slide_layout)
Image3_slide_title = Image3_slide.placeholders[0]
pic3 = Image3_slide.placeholders[1].insert_picture('image3.jpg')
Image3_slide_caption = Image3_slide.placeholders[2]

Image3_slide_title.text = "Simple Work Place"
Image3_slide_caption.text = "Nice Background!"

# IMAGE4 SLIDE
Image4_slide = prs.slides.add_slide(image_slide_layout)
Image4_slide_title = Image4_slide.placeholders[0]
pic4 = Image4_slide.placeholders[1].insert_picture('image4.jpg')
Image4_slide_caption = Image4_slide.placeholders[2]

Image4_slide_title.text = "Modern Furniture"
Image4_slide_caption.text = "Too Expensive"

# IMAGE5 SLIDE
Image5_slide = prs.slides.add_slide(image_slide_layout)
Image5_slide_title = Image5_slide.placeholders[0]
pic4 = Image5_slide.placeholders[1].insert_picture('image5.jpg')
Image5_slide_caption = Image5_slide.placeholders[2]

Image5_slide_title.text = "Cameraaa"
Image5_slide_caption.text = "Analog one"

"""

# SAVE THE PRESENTATION
prs.save('presentation.pptx')
